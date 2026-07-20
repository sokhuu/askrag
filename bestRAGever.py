import os
import json
import base64
import glob
import hashlib
import pickle
import re
import threading
from collections import defaultdict
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import List, Optional

import openai
from dotenv import load_dotenv
from pydantic import BaseModel
from tenacity import retry, retry_if_exception_type, stop_after_attempt, wait_random_exponential

from langchain_core.documents import Document
from langchain_core.messages import HumanMessage, SystemMessage, AIMessage
from langchain_openai import ChatOpenAI, OpenAIEmbeddings
from langchain_chroma import Chroma
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.retrievers import BM25Retriever
from langchain_classic.retrievers.ensemble import EnsembleRetriever
from langchain_cohere import CohereRerank

load_dotenv()

DOCS_PATH = "docs"
PERSIST_DIRECTORY = "db/chroma_combined"
ANSWER_CACHE_PATH = f"{PERSIST_DIRECTORY}.answer_cache.json"
EVENT_CACHE_PATH = f"{PERSIST_DIRECTORY}.event_cache.json"
BM25_CACHE_PATH = f"{PERSIST_DIRECTORY}.bm25.pkl"
MAX_HISTORY_TURNS = 10
TEXT_EXTENSIONS = {".txt"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg"}
PDF_EXTENSIONS = {".pdf"}
DOCX_EXTENSIONS = {".docx"}
PPTX_EXTENSIONS = {".pptx"}
SPREADSHEET_EXTENSIONS = {".xlsx", ".csv"}

PDF_STRATEGY = os.getenv("PDF_STRATEGY", "fast")  # override with "hi_res" for table/image extraction

embedding_model = OpenAIEmbeddings(model="text-embedding-3-small")
llm = ChatOpenAI(model="gpt-4o", temperature=0)
llm_mini = ChatOpenAI(model="gpt-4o-mini", temperature=0)


@retry(
    wait=wait_random_exponential(min=1, max=60),
    stop=stop_after_attempt(6),
    retry=retry_if_exception_type(openai.RateLimitError),
    reraise=True,
)
def _invoke_with_retry(model, messages):
    return model.invoke(messages)


class QueryVariations(BaseModel):
    queries: List[str]


class EventIntentResult(BaseModel):
    event_id: Optional[str] = None


# ingestion

def partition_pdf_file(file_path, strategy=None):
    """Extract chunks from a PDF via unstructured. strategy='fast' (default) for text-only;
    'hi_res' enables table structure and image extraction at the cost of much slower parsing."""
    from unstructured.partition.pdf import partition_pdf
    from unstructured.chunking.title import chunk_by_title

    strategy = strategy or PDF_STRATEGY
    hi_res = strategy == "hi_res"
    elements = partition_pdf(
        filename=file_path,
        strategy=strategy,
        infer_table_structure=hi_res,
        extract_image_block_types=["Image"] if hi_res else [],
        extract_image_block_to_payload=hi_res,
    )
    chunks = chunk_by_title(
        elements,
        max_characters=3000,
        new_after_n_chars=2400,
        combine_text_under_n_chars=500,
    )

    normalized = []
    for chunk in chunks:
        tables, images = [], []
        if hasattr(chunk, "metadata") and hasattr(chunk.metadata, "orig_elements"):
            for element in chunk.metadata.orig_elements:
                element_type = type(element).__name__
                if element_type == "Table":
                    tables.append(getattr(element.metadata, "text_as_html", element.text))
                elif element_type == "Image" and hasattr(element.metadata, "image_base64"):
                    images.append(element.metadata.image_base64)
        normalized.append({"text": chunk.text, "tables": tables, "images": images, "source": file_path})
    return normalized


def partition_text_file(file_path, chunk_size=1000, chunk_overlap=200):
    """Split a plain text file into text-only chunks"""
    with open(file_path, "r", encoding="utf-8") as f:
        raw_text = f.read()

    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    pieces = splitter.split_text(raw_text)
    return [{"text": piece, "tables": [], "images": [], "source": file_path} for piece in pieces]


def partition_image_file(file_path):
    """Treat a standalone image file as a single image-only chunk"""
    with open(file_path, "rb") as f:
        image_base64 = base64.b64encode(f.read()).decode("utf-8")
    return [{"text": "", "tables": [], "images": [image_base64], "source": file_path}]


def partition_docx_file(file_path, chunk_size=1000, chunk_overlap=200):
    from docx import Document as DocxDocument
    doc = DocxDocument(file_path)
    full_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    chunks = [
        {"text": piece, "tables": [], "images": [], "source": file_path}
        for piece in splitter.split_text(full_text)
    ]
    for table in doc.tables:
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        html = "<table>" + "".join(
            "<tr>" + "".join(f"<td>{c}</td>" for c in row) + "</tr>" for row in rows
        ) + "</table>"
        text = "\n".join("\t".join(row) for row in rows)
        chunks.append({"text": text, "tables": [html], "images": [], "source": file_path})
    return chunks or [{"text": "", "tables": [], "images": [], "source": file_path}]


def partition_pptx_file(file_path, chunk_size=1000, chunk_overlap=200):
    from pptx import Presentation
    prs = Presentation(file_path)
    slide_texts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [s.text.strip() for s in slide.shapes if hasattr(s, "text") and s.text.strip()]
        if texts:
            slide_texts.append(f"Slide {i}:\n" + "\n".join(texts))
    full_text = "\n\n".join(slide_texts)
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    pieces = splitter.split_text(full_text)
    return [{"text": piece, "tables": [], "images": [], "source": file_path} for piece in pieces] \
        or [{"text": "", "tables": [], "images": [], "source": file_path}]


def partition_spreadsheet_file(file_path, rows_per_chunk=100):
    import pandas as pd
    ext = os.path.splitext(file_path)[1].lower()
    sheets = (
        [("", pd.read_csv(file_path))]
        if ext == ".csv"
        else [(name, pd.read_excel(file_path, sheet_name=name))
              for name in pd.ExcelFile(file_path).sheet_names]
    )
    chunks = []
    for sheet_name, df in sheets:
        for start in range(0, max(len(df), 1), rows_per_chunk):
            sl = df.iloc[start:start + rows_per_chunk]
            label = f"Sheet '{sheet_name}', " if sheet_name else ""
            text = f"{label}rows {start}–{start + len(sl) - 1}\n{sl.to_string(index=False)}"
            chunks.append({"text": text, "tables": [sl.to_html(index=False)], "images": [], "source": file_path})
    return chunks


def create_ai_enhanced_summary(text, tables, images, source, model=None):
    """Create a rich, searchable description for mixed content.
    Uses gpt-4o (vision) by default; pass llm_mini for text/table-only chunks."""
    model = model or llm
    try:
        prompt_text = f"""You are creating a searchable description for document content retrieval.

Source file: {source}

TEXT CONTENT:
{text}

"""
        if tables:
            prompt_text += "TABLES:\n"
            for i, table in enumerate(tables):
                prompt_text += f"Table {i+1}:\n{table}\n\n"

        prompt_text += """
YOUR TASK:
Generate a comprehensive, searchable description that covers:
1. Key facts, numbers, and data points from text and tables
2. Main topics and concepts discussed
3. Questions this content could answer
4. Visual content analysis (charts, diagrams, patterns in images)
5. Alternative search terms users might use

Make it detailed and searchable - prioritize findability over brevity.

SEARCHABLE DESCRIPTION:"""

        message_content = [{"type": "text", "text": prompt_text}]
        for image_base64 in images:
            message_content.append({
                "type": "image_url",
                "image_url": {"url": f"data:image/jpeg;base64,{image_base64}"},
            })

        response = _invoke_with_retry(model, [HumanMessage(content=message_content)])
        return response.content
    except Exception as e:
        print(f"  AI summary failed for {source}: {e}")
        summary = f"{text[:300]}..." if text else f"[Image from {source}]"
        if tables:
            summary += f" [Contains {len(tables)} table(s)]"
        if images:
            summary += f" [Contains {len(images)} image(s)]"
        return summary


def _save_images_to_disk(images_base64, images_dir):
    """Write base64 image blobs to disk keyed by content hash. Returns list of paths."""
    os.makedirs(images_dir, exist_ok=True)
    paths = []
    for img_b64 in images_base64:
        img_hash = hashlib.sha256(img_b64.encode()).hexdigest()[:16]
        img_path = os.path.join(images_dir, f"{img_hash}.jpg")
        if not os.path.exists(img_path):
            with open(img_path, "wb") as f:
                f.write(base64.b64decode(img_b64))
        paths.append(img_path)
    return paths


def _summarize_one(index, total, chunk, images_dir):
    """Summarize a single chunk; returns (original_index, Document) for stable ordering."""
    print(f"Summarizing chunk {index}/{total} ({chunk['source']})...")
    if chunk["images"]:
        # Vision content — gpt-4o required
        searchable_content = create_ai_enhanced_summary(
            chunk["text"], chunk["tables"], chunk["images"], chunk["source"]
        )
    elif chunk["tables"]:
        # Text + tables, no images — gpt-4o-mini is sufficient (~10x cheaper)
        searchable_content = create_ai_enhanced_summary(
            chunk["text"], chunk["tables"], [], chunk["source"], model=llm_mini
        )
    else:
        searchable_content = chunk["text"]

    image_paths = _save_images_to_disk(chunk["images"], images_dir)

    return index, Document(
        page_content=searchable_content,
        metadata={
            "source": chunk["source"],
            "filename": os.path.basename(chunk["source"]),
            "file_type": os.path.splitext(chunk["source"])[1].lstrip(".").lower(),
            "has_tables": bool(chunk["tables"]),
            "has_images": bool(chunk["images"]),
            "image_paths": json.dumps(image_paths),
            "original_content": json.dumps({
                "raw_text": chunk["text"],
                "tables_html": chunk["tables"],
            }),
        },
    )

_print_lock = threading.Lock()

def summarize_chunks(raw_chunks, images_dir, max_workers=20):
    """Turn normalized chunks into LangChain Documents, parallelizing GPT-4o calls."""
    total = len(raw_chunks)
    results = [None] * total

    with ThreadPoolExecutor(max_workers=max_workers) as pool:
        futures = {
            pool.submit(_summarize_one, i, total, chunk, images_dir): (i, chunk["source"])
            for i, chunk in enumerate(raw_chunks)
        }
        done_count = [0]
        for future in as_completed(futures):
            idx, source = futures[future]
            try:
                _, doc = future.result()
                results[idx] = doc
            except Exception as e:
                with _print_lock:
                    print(f"\n  [SKIP] Summarization failed for chunk {idx} ({source}): {e}")
            with _print_lock:
                done_count[0] += 1
                print(f"  [{done_count[0]}/{total}] chunks summarized", end="\r", flush=True)

    print()
    return [r for r in results if r is not None]


# vector store


def manifest_path(persist_directory):
    return f"{persist_directory}.manifest.json"


def _hash_file(file_path):
    """Stream-hash a file in 64 KB blocks to avoid loading large PDFs into memory."""
    h = hashlib.sha256()
    with open(file_path, "rb") as f:
        for block in iter(lambda: f.read(65536), b""):
            h.update(block)
    return h.hexdigest()


def compute_per_file_signatures(docs_path):
    """Return {file_path: sha256_hex} for every supported file under docs_path."""
    all_extensions = TEXT_EXTENSIONS | IMAGE_EXTENSIONS | PDF_EXTENSIONS | DOCX_EXTENSIONS | PPTX_EXTENSIONS | SPREADSHEET_EXTENSIONS
    signatures = {}
    for file_path in sorted(glob.glob(os.path.join(docs_path, "**", "*"), recursive=True)):
        if os.path.isfile(file_path) and os.path.splitext(file_path)[1].lower() in all_extensions:
            signatures[file_path] = _hash_file(file_path)
    return signatures


def read_manifest(persist_directory):
    """Return stored per-file signature dict. Old single-hash format treated as empty (triggers rebuild)."""
    path = manifest_path(persist_directory)
    if not os.path.exists(path):
        return {}
    with open(path, "r") as f:
        data = json.load(f)
    return data.get("files", {})


def write_manifest(persist_directory, file_signatures):
    with open(manifest_path(persist_directory), "w") as f:
        json.dump({"files": file_signatures}, f, indent=2)


def _partition_one(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    with _print_lock:
        print(f"Partitioning {file_path}...")
    try:
        if ext in TEXT_EXTENSIONS:
            return partition_text_file(file_path)
        elif ext in PDF_EXTENSIONS:
            return partition_pdf_file(file_path)
        elif ext in IMAGE_EXTENSIONS:
            return partition_image_file(file_path)
        elif ext in DOCX_EXTENSIONS:
            return partition_docx_file(file_path)
        elif ext in PPTX_EXTENSIONS:
            return partition_pptx_file(file_path)
        elif ext in SPREADSHEET_EXTENSIONS:
            return partition_spreadsheet_file(file_path)
        return []
    except Exception as e:
        with _print_lock:
            print(f"  [SKIP] Failed to partition {file_path}: {e}")
        return []


def partition_files(file_paths, max_workers=4):
    """Partition a list of files into normalized chunks in parallel."""
    all_chunks = []
    with ThreadPoolExecutor(max_workers=max_workers) as pool:
        futures = {pool.submit(_partition_one, fp): fp for fp in sorted(file_paths)}
        for future in as_completed(futures):
            all_chunks.extend(future.result())
    return all_chunks


def build_or_load_vectorstore(docs_path, persist_directory):
    current_signatures = compute_per_file_signatures(docs_path)
    stored_signatures = read_manifest(persist_directory)

    added = [p for p in current_signatures if p not in stored_signatures]
    modified = [p for p in current_signatures if p in stored_signatures and current_signatures[p] != stored_signatures[p]]
    deleted = [p for p in stored_signatures if p not in current_signatures]

    if os.path.exists(persist_directory) and not added and not modified and not deleted:
        print("Vector store is up to date.")
        return Chroma(
            persist_directory=persist_directory,
            embedding_function=embedding_model,
            collection_metadata={"hnsw:space": "cosine"},
        )

    if added or modified or deleted:
        summary = []
        if added:
            summary.append(f"{len(added)} added")
        if modified:
            summary.append(f"{len(modified)} modified")
        if deleted:
            summary.append(f"{len(deleted)} deleted")
        print(f"Detected changes: {', '.join(summary)}. Running incremental update...")
    else:
        print("No existing vector store found. Running full ingestion pipeline...")

    vectorstore = Chroma(
        persist_directory=persist_directory,
        embedding_function=embedding_model,
        collection_metadata={"hnsw:space": "cosine"},
    )

    # Remove chunks for deleted or modified files
    stale_files = deleted + modified
    if stale_files:
        print(f"Removing stale chunks for {len(stale_files)} file(s)...")
        for file_path in stale_files:
            result = vectorstore.get(where={"source": file_path})
            if result["ids"]:
                vectorstore.delete(ids=result["ids"])
                print(f"  Removed {len(result['ids'])} chunk(s) from {file_path}")

    # Ingest added and modified files
    images_dir = os.path.join(persist_directory, "images")
    to_ingest = added + modified
    if to_ingest:
        print(f"Ingesting {len(to_ingest)} file(s)...")
        raw_chunks = partition_files(to_ingest)
        if raw_chunks:
            documents = summarize_chunks(raw_chunks, images_dir=images_dir)
            vectorstore.add_documents(documents)
            print(f"Added {len(documents)} chunk(s) to vector store.")
    elif not os.path.exists(persist_directory):
        raise FileNotFoundError(f"No supported documents found in {docs_path}.")

    # Docs changed — invalidate cached BM25 index and answer cache
    for stale_path in [BM25_CACHE_PATH, ANSWER_CACHE_PATH]:
        if os.path.exists(stale_path):
            os.remove(stale_path)

    write_manifest(persist_directory, current_signatures)
    print(f"Vector store ready. Tracking {len(current_signatures)} file(s).")
    return vectorstore


def _load_all_documents(vectorstore):
    raw = vectorstore.get(include=["documents", "metadatas"])
    return [
        Document(page_content=content, metadata=metadata)
        for content, metadata in zip(raw["documents"], raw["metadatas"])
    ]


def _build_or_load_bm25_retriever(vectorstore, k):
    """Load a pickled BM25 index from disk, or build and cache it on first run."""
    index_path = BM25_CACHE_PATH
    if os.path.exists(index_path):
        print("Loading BM25 index from disk...")
        with open(index_path, "rb") as f:
            retriever = pickle.load(f)
        retriever.k = k
        return retriever

    print("Building BM25 index (first run or docs changed)...")
    documents = _load_all_documents(vectorstore)
    retriever = BM25Retriever.from_documents(documents)
    retriever.k = k
    with open(index_path, "wb") as f:
        pickle.dump(retriever, f)
    return retriever


def build_hybrid_retriever(vectorstore, k=10):
    vector_retriever = vectorstore.as_retriever(search_kwargs={"k": k})
    bm25_retriever = _build_or_load_bm25_retriever(vectorstore, k)
    return EnsembleRetriever(retrievers=[vector_retriever, bm25_retriever], weights=[0.7, 0.3])


# accuracy pipeline for history-aware rewrite -> multi-query -> RRF -> rerank

def rewrite_standalone_query(user_question, chat_history):
    if not chat_history:
        return user_question
    messages = [
        SystemMessage(content=(
            "Given the chat history, rewrite the new question to be a standalone, "
            "searchable query. Just return the rewritten question."
        )),
    ] + chat_history + [HumanMessage(content=f"New question: {user_question}")]
    result = llm.invoke(messages)
    return result.content.strip()


def generate_query_variations(standalone_query, n=3):
    llm_with_structure = llm.with_structured_output(QueryVariations)
    prompt = f"""Generate {n} different variations of this query that would help retrieve relevant documents:

Original query: {standalone_query}

Return {n} alternative queries that rephrase or approach the same question from different angles."""
    response = llm_with_structure.invoke(prompt)
    return response.queries


def reciprocal_rank_fusion(chunk_lists, k=60):
    rrf_scores = defaultdict(float)
    unique_chunks = {}
    for chunks in chunk_lists:
        for position, chunk in enumerate(chunks, 1):
            content = chunk.page_content
            unique_chunks[content] = chunk
            rrf_scores[content] += 1 / (k+position)

    return sorted(
        ((unique_chunks[content], score) for content, score in rrf_scores.items()),
        key=lambda x: x[1],
        reverse=True,
    )

_JUNK_TEXT_PATTERN = re.compile(
    r"learn more|finra\.org/rules-guidance|\bUP\b|VERSIONS\b|"
    r"\d{1,2}/\d{1,2}/\d{2,4},\s*\d{1,2}:\d{2}\s*[AP]M|[‹›]",
    re.I,
)


def _is_junk_document(doc):
    """Flag PDF navigation/footer chunks (breadcrumbs, prev/next links, version-history
    links, page URLs) that were captured as literal text when the FINRA rule pages were
    saved to PDF. These can otherwise outrank real rule text via exact keyword matches
    on the rule number/title that appears in the nav chrome."""
    original = json.loads(doc.metadata.get("original_content", "{}"))
    text = original.get("raw_text", "") or doc.page_content
    stripped = _JUNK_TEXT_PATTERN.sub(" ", text)
    stripped = re.sub(r"https?://\S+", " ", stripped)
    stripped = re.sub(r"\s+", " ", stripped).strip()
    return len(stripped) < 120


def filter_junk_documents(docs):
    filtered = [d for d in docs if not _is_junk_document(d)]
    return filtered or docs


_reranker = None
_rerank_unavailable = False

def get_reranker(top_n):
    """Build the Cohere reranker once; fall back permanently if no CO_API_KEY is set"""
    global _reranker, _rerank_unavailable
    if _rerank_unavailable:
        return None
    if _reranker is None:
        try:
            _reranker = CohereRerank(model="rerank-english-v3.0", top_n=top_n)
        except Exception as e:
            print(f"Cohere reranker unavailable ({e}). Falling back to RRF ranking only.")
            _rerank_unavailable = True
            return None
    return _reranker


def rerank_candidates(candidates, query, top_n=5):
    reranker = get_reranker(top_n)
    if reranker is None:
        return candidates[:top_n]
    try:
        return reranker.compress_documents(candidates, query)
    except Exception as e:
        print(f"Reranking failed ({e}). Falling back to RRF ranking only.")
        return candidates[:top_n]


# answer generation


def generate_answer(top_docs, user_question, chat_history):
    prompt_text = f"""Based on the following documents from multiple sources, answer this question directly and concisely: {user_question}

If the answer requires combining information from multiple documents or subjects, synthesize them clearly. If the documents don't contain enough information, say "I don't have enough information to answer that question based on the provided documents."

CONTENT TO ANALYZE:
"""
    image_blocks = []
    for i, doc in enumerate(top_docs, 1):
        prompt_text += f"\n--- Document {i} (source: {doc.metadata.get('source', 'unknown')}) ---\n"
        original = json.loads(doc.metadata.get("original_content", "{}"))

        raw_text = original.get("raw_text", "")
        if raw_text:
            prompt_text += f"TEXT:\n{raw_text}\n"

        for j, table in enumerate(original.get("tables_html", []), 1):
            prompt_text += f"TABLE {j}:\n{table}\n"

        for img_path in json.loads(doc.metadata.get("image_paths", "[]")):
            if os.path.exists(img_path):
                with open(img_path, "rb") as f:
                    image_base64 = base64.b64encode(f.read()).decode("utf-8")
                image_blocks.append({
                    "type": "image_url",
                    "image_url": {"url": f"data:image/jpeg;base64,{image_base64}"},
                })

    prompt_text += "\nANSWER:"
    message_content = [{"type": "text", "text": prompt_text}] + image_blocks

    messages = [
        SystemMessage(content=(
            "You are a direct, helpful assistant that answers using only the provided "
            "documents and conversation history."
        )),
    ] + chat_history + [HumanMessage(content=message_content)]

    print()
    full_response = ""
    for chunk in llm.stream(messages):
        print(chunk.content, end="", flush=True)
        full_response += chunk.content
    print()
    return full_response


# event-triggered obligations checklists

EVENT_TEMPLATES = {
    "onboard_rep": {
        "label": "Onboard a registered representative",
        "hint": "New hire who will register as a rep",
        "topics": [
            {"category": "Registration & Form U4", "query": "What FINRA registration requirements and Form U4 filing obligations apply when a firm hires a new registered representative?"},
            {"category": "Fingerprinting", "query": "What fingerprinting requirements apply to a newly hired associated person under FINRA rules?"},
            {"category": "Qualification Exams", "query": "What qualification examinations must a new registered representative pass before conducting securities business?"},
            {"category": "Continuing Education", "query": "What continuing education requirements and deadlines, including the annual Regulatory Element completion date, apply to a newly registered person?"},
        ],
    },
    "terminate_rep": {
        "label": "Terminate a registered representative",
        "hint": "Rep leaving the firm, voluntary or involuntary",
        "topics": [
            {"category": "Form U5 Filing", "query": "What is the deadline and process for filing a Form U5 when a registered representative is terminated?"},
            {"category": "Outstanding Complaints & Disputes", "query": "What obligations does a firm have regarding a terminated registered person's outstanding customer complaints or disputes?"},
        ],
    },
    "open_branch": {
        "label": "Open a new branch office",
        "hint": "Establishing a new branch or OSJ location",
        "topics": [
            {"category": "Branch Supervision", "query": "What supervision requirements apply to a newly opened branch office under FINRA Rule 3110?"},
            {"category": "Inspection Cycle", "query": "What are the inspection cycle requirements for a new branch office?"},
            {"category": "Branch Registration", "query": "What registration or notice filing requirements apply when opening a new branch office?"},
        ],
    },
    "aum_headcount_threshold": {
        "label": "Cross an AUM or headcount threshold",
        "hint": "Firm growth trips new regulatory thresholds",
        "topics": [
            {"category": "Size-Based Obligations", "query": "What new FINRA obligations apply to a member firm when it crosses significant size or capital thresholds?"},
            {"category": "Supervisory Controls", "query": "What supervisory control system requirements apply as a firm grows under FINRA Rule 3120?"},
            {"category": "Restricted Firm Obligations", "query": "What restricted firm obligations apply based on firm size or disciplinary history under FINRA Rule 4111?"},
        ],
    },
    "customer_complaint": {
        "label": "Receive a customer complaint",
        "hint": "A written complaint just came in",
        "topics": [
            {"category": "Complaint Reporting", "query": "What are a member firm's reporting requirements when it receives a written customer complaint under FINRA Rule 4530?"},
            {"category": "Complaint Recordkeeping", "query": "What recordkeeping requirements apply to written customer complaints under FINRA Rule 4513?"},
        ],
    },
    "outside_business_activity": {
        "label": "Rep wants an outside business activity",
        "hint": "Outside business activity or private securities transaction",
        "topics": [
            {"category": "OBA Disclosure & Approval", "query": "What disclosure and approval obligations apply when a registered person wants to engage in an outside business activity under FINRA Rule 3270?"},
            {"category": "Private Securities Transactions", "query": "What obligations apply to a registered person's private securities transactions under FINRA Rule 3280?"},
        ],
    },
}


def classify_event_intent(question):
    """Detect whether a free-form question is actually describing one of the canned
    business events happening (e.g. "we just hired a new rep") rather than a general
    question about a rule. Returns a matching event_id, or None."""
    event_list = "\n".join(f"- {eid}: {t['label']} ({t['hint']})" for eid, t in EVENT_TEMPLATES.items())
    prompt = f"""A user of a FINRA compliance tool wrote the message below. Decide whether it describes one of these business events just happening, or about to happen, at their firm (not a general question about what a rule says):

{event_list}

Message: "{question}"

Return the matching event_id only if the message clearly describes one of these events (e.g. "we just hired a new rep", "I'm terminating someone", "we're opening a new office", "a client just complained"). Return null if it's a general question or doesn't clearly match any event."""
    llm_with_structure = llm_mini.with_structured_output(EventIntentResult)
    result = llm_with_structure.invoke(prompt)
    return result.event_id if result.event_id in EVENT_TEMPLATES else None


def _load_event_cache():
    if not os.path.exists(EVENT_CACHE_PATH):
        return {}
    with open(EVENT_CACHE_PATH) as f:
        return json.load(f)


def _save_event_cache(cache):
    with open(EVENT_CACHE_PATH, "w") as f:
        json.dump(cache, f, indent=2)


class ChecklistItemModel(BaseModel):
    category: str
    obligation: str
    rule: str
    deadline: str
    due_days: Optional[int] = None


class ChecklistResult(BaseModel):
    items: List[ChecklistItemModel]


def generate_checklist_items(top_docs, event_label, context, topics):
    topics_text = "\n".join(f"- {t['category']}: {t['query']}" for t in topics)
    prompt_text = f"""A brokerage compliance officer just experienced this event: "{event_label}".
{f'Additional context they provided: {context}' if context else ''}

Produce a structured obligations checklist so nothing falls through the cracks, covering exactly these obligation areas (use the short label before the colon verbatim as the "category" for its items — do not invent new categories, and do not add a catch-all/"unspecified" item):
{topics_text}

For each obligation, cite the specific FINRA rule number and name from the source documents (e.g. "FINRA Rule 3110 (Supervision)").

Deadline accuracy is critical — these are read by compliance officers who will rely on them. A "deadline" includes any timing or triggering condition the rule attaches to the obligation, not just a day-count — e.g. "before registration becomes effective", "annually by December 31", "promptly", and "within 30 days" are all valid deadlines to report; report them using that same wording. Follow these rules strictly:
- Only state a deadline if the source text ties it to THIS specific obligation. Do not borrow or infer a deadline from a nearby but different provision in the same rule (for example, a recordkeeping/"available promptly upon regulatory request" clause is NOT the deadline for the filing action itself — those are different obligations even when they appear in the same paragraph, even if no other timing language exists for the filing action).
- Quote or closely paraphrase the deadline exactly as the rule states it, including any specific date (e.g. "annually by December 31", not just "annually") or condition (e.g. "before registration becomes effective", not just "eventually").
- Only if the source text truly states no timing or triggering condition at all for this specific obligation, set deadline to "No specific deadline stated in the rule text" and due_days to null. Do not invent a plausible-sounding deadline, but do not discard a real one either.
- Only set due_days to a non-null integer if the rule states a fixed number of days from a clearly identifiable trigger event for THIS obligation (e.g. "within 30 days of Form U4 filing"). For conditional deadlines without a day-count (e.g. "before registration becomes effective", "annually by December 31"), leave due_days null but still report the condition in the deadline field.

If the provided documents don't have enough information for one of the areas above, add a single item under that same category noting the gap, with deadline "Not covered by the retrieved rule text" and due_days null.

CONTENT TO ANALYZE:
"""
    for i, doc in enumerate(top_docs, 1):
        prompt_text += f"\n--- Document {i} (source: {doc.metadata.get('source', 'unknown')}) ---\n"
        original = json.loads(doc.metadata.get("original_content", "{}"))
        raw_text = original.get("raw_text", "")
        if raw_text:
            prompt_text += f"TEXT:\n{raw_text}\n"
        for j, table in enumerate(original.get("tables_html", []), 1):
            prompt_text += f"TABLE {j}:\n{table}\n"

    llm_with_structure = llm.with_structured_output(ChecklistResult)
    result = llm_with_structure.invoke(prompt_text)
    return [item.model_dump() for item in result.items]


def _format_checklist_text(items):
    by_category = {}
    for item in items:
        by_category.setdefault(item["category"], []).append(item)

    lines = []
    for category, cat_items in by_category.items():
        lines.append(f"### {category}")
        for item in cat_items:
            lines.append(f"- {item['obligation']} — {item['rule']} ({item['deadline']})")
        lines.append("")
    return "\n".join(lines).strip()


def answer_event(event_id, context, hybrid_retriever, per_topic_fetch_k=10, per_topic_final_k=4):
    template = EVENT_TEMPLATES.get(event_id)
    if template is None:
        raise ValueError(f"Unknown event_id: {event_id}")

    context = (context or "").strip()
    cache_key = hashlib.sha256(f"{event_id}|{context}".lower().encode()).hexdigest()[:16]
    cache = _load_event_cache()
    if cache_key in cache:
        items = cache[cache_key]
        return template["label"], _format_checklist_text(items), items

    # Retrieve and rerank per topic (rather than pooling all topics into one shared
    # top-k) so every obligation area gets guaranteed dedicated coverage — otherwise
    # a topic's supporting chunk can be crowded out of a shared pool by unrelated
    # topics, and that's fragile to reranker run-to-run variance on borderline scores.
    def _top_docs_for_topic(topic):
        query = f"{topic['query']} {context}".strip()
        results = hybrid_retriever.invoke(query)
        results = [d for d in results if not _is_junk_document(d)] or results
        return rerank_candidates(results[:per_topic_fetch_k], query, top_n=per_topic_final_k)

    with ThreadPoolExecutor() as pool:
        per_topic_docs = list(pool.map(_top_docs_for_topic, template["topics"]))

    top_docs = []
    seen = set()
    for docs in per_topic_docs:
        for doc in docs:
            if doc.page_content not in seen:
                seen.add(doc.page_content)
                top_docs.append(doc)

    items = generate_checklist_items(top_docs, template["label"], context, template["topics"])

    cache[cache_key] = items
    _save_event_cache(cache)

    return template["label"], _format_checklist_text(items), items


# CHAT :3

def _load_answer_cache():
    if not os.path.exists(ANSWER_CACHE_PATH):
        return {}
    with open(ANSWER_CACHE_PATH) as f:
        return json.load(f)


def _save_answer_cache(cache):
    with open(ANSWER_CACHE_PATH, "w") as f:
        json.dump(cache, f, indent=2)


def ask_question(user_question, hybrid_retriever, chat_history, fetch_k=15, final_k=5,
                 vectorstore=None, metadata_filter=None):
    standalone_query = rewrite_standalone_query(user_question, chat_history)

    cache_key = hashlib.sha256(standalone_query.lower().strip().encode()).hexdigest()[:16]
    cache = _load_answer_cache()
    if cache_key in cache:
        print(f"\n{cache[cache_key]}")
        chat_history.append(HumanMessage(content=user_question))
        chat_history.append(AIMessage(content=cache[cache_key]))
        return cache[cache_key]

    variations = generate_query_variations(standalone_query)
    all_queries = [standalone_query] + variations

    if metadata_filter and vectorstore:
        retriever = vectorstore.as_retriever(search_kwargs={"k": fetch_k, "filter": metadata_filter})
        with ThreadPoolExecutor() as pool:
            all_results = list(pool.map(retriever.invoke, all_queries))
    else:
        with ThreadPoolExecutor() as pool:
            all_results = list(pool.map(hybrid_retriever.invoke, all_queries))

    fused = reciprocal_rank_fusion(all_results, k=60)
    fused = [pair for pair in fused if not _is_junk_document(pair[0])] or fused
    candidates = [doc for doc, _ in fused[:fetch_k]]

    top_docs = rerank_candidates(candidates, standalone_query, top_n=final_k)

    answer = generate_answer(top_docs, user_question, chat_history)

    cache[cache_key] = answer
    _save_answer_cache(cache)

    chat_history.append(HumanMessage(content=user_question))
    chat_history.append(AIMessage(content=answer))

    # Keep history bounded to avoid context overflow
    if len(chat_history) > MAX_HISTORY_TURNS * 2:
        chat_history[:] = chat_history[-(MAX_HISTORY_TURNS * 2):]

    return answer


def start_chat():
    vectorstore = build_or_load_vectorstore(DOCS_PATH, PERSIST_DIRECTORY)
    hybrid_retriever = build_hybrid_retriever(vectorstore)

    chat_history = []
    print("\nAsk me about FINRA compliance.")
    while True:
        question = input("\nYour question: ")
        if question.lower() == "quit":
            print("Goodbye.")
            break
        ask_question(question, hybrid_retriever, chat_history, vectorstore=vectorstore)


if __name__ == "__main__":
    start_chat()
